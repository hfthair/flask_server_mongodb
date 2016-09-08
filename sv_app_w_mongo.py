#!/usr/bin/python
# -*- coding: utf-8 -*-

from flask import (Flask, request, redirect, Response, url_for, 
                    send_from_directory, jsonify, send_file, 
                    render_template, flash
                    )

from flask_pymongo import PyMongo
from gridfs import GridFS
from bson.objectid import ObjectId
from bson import json_util

# from collections import OrderedDict
import os, uuid, json, time, shutil
import traceback


app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 25<<20  # max upload size < 16M
app.config['MONGO_DBNAME'] = 'sv4app'
app.debug = False

mongo = PyMongo(app)
ctx = app.test_request_context('/')
ctx.push()
db = mongo.db
fs = GridFS(db)
db.usr.create_index('usr', unique=True)
db.meeting.create_index('name')
ctx.pop()

MEETING_MAP = {}
ALLOWED_EXT = ['.pdf', '.doc', '.ppt', '.pptx', '.docx', '.xls', '.xlsx', '.caj', '.txt', '.html', '.rec', '.commonts']


#############################################################################################################
############################################### login #######################################################

import base64
from flask_login import (LoginManager, current_user, login_required,
                            login_user, logout_user, UserMixin,
                            confirm_login, fresh_login_required)

app.config["SECRET_KEY"] = "ITSASnrmUTYewqbECRETkutepVC.XddJGrio3fKHFKDkdkls-KDSkdkfpLFDOSF"
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login_failed'

class User(UserMixin):
    def __init__(self, id, name, passwd, active=True):
        self.id = id
        self.name = name
        self.passwd = passwd
        self.active = active

    def is_active(self):
        return self.active

    @staticmethod
    def get_from_id(id):
        idobj = ObjectId(id)
        res = db.usr.find_one(dict(_id=idobj))
        if res != None:
            name = res.get('usr', '')
            passwd = res.get('passwd', '')
            isactive = res.get('active', True)
            return User(id, name, passwd, isactive)


    @staticmethod
    def get_from_name(name):
        res = db.usr.find_one(dict(usr=name))
        if res == None:
            return None
        else:
            id = res.get('_id', '')
            passwd = res.get('passwd', '')
            isactive = res.get('active', True)
            if id != '' and passwd != '':
                #id.generation_time
                return User(str(id), name, passwd, isactive)
            

    @staticmethod
    def regist_new_usr(_usr, _passwd):
        if _usr == None or _usr == '' or _passwd == None or _passwd == '':
            return False
        temp = dict(usr = _usr, passwd = _passwd, active = True)
        try:
            db.usr.insert_one(temp)
        except:
            traceback.print_exc()
            return False
        else:
            return True
    
    #test port
    @staticmethod
    def get_all_usr_name():
        dr = db.usr.find()
        res = []
        for i in dr:
            res.append(i.get('usr', ''))
        return res



@login_manager.user_loader
def load_user(id):
    return User.get_from_id(id)

@login_manager.request_loader
def load_user(request):
    api_key = request.args.get('api_key')
    if api_key:
        user, passwd = api_key.split(':')
        cUser = User.get_from_name(user)
        if cUser.name == user and cUser.passwd == passwd and cUser.is_active():
            return cUser
    api_key = request.headers.get('Authorization')
    if api_key:
        api_key = api_key.replace('Basic ', '', 1)
        try:
            api_key = base64.b64decode(api_key)
        except TypeError:
            traceback.print_exc()
        user, passwd = api_key.split(':')
        cUser = User.get_from_name(user)
        if cUser.name == user and cUser.passwd == passwd and cUser.is_active():
            return cUser
    return None

#############################################################################################################
import pptx
def pptx_to_desc(f):
    text_runs = []
    try:
        prtn = pptx.Presentation(f)
        for slide in prtn.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    except:
        traceback.print_exc()
    return '\n'.join(text_runs)

def datetime_to_timetimestr(t):
    return str(time.mktime(t.timetuple()) + t.microsecond / 1E6)

def file_ext_check(name):
    return os.path.splitext(name)[1] in ALLOWED_EXT

def secure_filename(s):
    return secure_name(s)

def secure_name(s):
    s = s.replace('\\', '_')
    s = s.replace('/', '_')
    s = s.replace('..', '_')
    s = s.replace(' ', '_')
    s = s.replace('(', '_')
    s = s.replace(')', '_')
    s = s.replace('[', '_')
    s = s.replace(']', '_')
    return s


def get_meeting_list(uid=''):
    all_meeting = {}
    try:
        res = db.meeting.find(limit=512)
        for i in res:
            desc = ''
            objid = i.get('_id')
            f1 = i.get('main_files')
            for k1 in f1:
                ar1 = f1.get(k1)
                if len(ar1) == 2:
                    desc += ar1[1]
                    desc += '\n'
            
            f2 = i.get('ref_files')
            for k2 in f2:
                ar2 = f2.get(k2)
                if len(ar2) == 2:
                    desc += ar2[1]
                    desc += '\n'

            ttt = i.get('desc')
            if ttt:
                desc = ttt
            id = str(objid)
            t = objid.generation_time 
            tf = datetime_to_timetimestr(t)
            all_meeting[id] = dict(
                title = i.get('title'),
                creator = i.get('creator'),
                time = tf,
                desc = desc
            )
    except:
        traceback.print_exc()
    return all_meeting

def create_uuid():
    # return uuid.uuid1().bytes
    return str(uuid.uuid1())

def op_meeting_new(meeting, creator):
    cfg = dict(
        title = meeting,
        creator = creator,
        main_files = {},
        ref_files = {},
        record = {},
        comments = {}
    )
    try:
        res = db.meeting.insert_one(cfg)
        id = str(res.inserted_id)
        return id
    except:
        traceback.print_exc()
    return None

def op_meeting_delete(meeting_id):
    try:
        res = db.meeting.delete_one(
            {'_id':ObjectId(meeting_id)}
            )
        return res.deleted_count
    except:
        traceback.print_exc()
    return 0


def op_meeting_find_by_id(meeting_id):
    try:
        objid = ObjectId(meeting_id)
        res = db.meeting.find_one({'_id':objid})
        return res
    except:
        traceback.print_exc()
    return None

def op_meeting_add_main_file(meeting_id, file, filename, author):
    try:
        objid = fs.put(file, filename=filename, dir='m')
        id = str(objid)
        ext = os.path.splitext(filename)[1]
        desc = ''
        if ext in ('.ppt', '.pptx'):
            desc = pptx_to_desc(file)

        res = db.meeting.update_one(
            {'_id':ObjectId(meeting_id)}, 
            {'$set':{
                "main_files."+id:[author, filename],
                "desc":desc
            }}
            )
        return res.modified_count
    except:
        traceback.print_exc()
    return 0


def op_meeting_add_ref_file(meeting_id, file, filename, author):
    try:
        objid = fs.put(file, filename=filename, dir='ref')
        id = str(objid)
        res = db.meeting.update_one(
            {'_id':ObjectId(meeting_id)}, 
            {'$set':{
                "ref_files."+id:[author, filename]
            }}
            )
        return res.modified_count
    except:
        traceback.print_exc()
    return 0

def op_meeting_rm_usr_file(meeting_id, file_id, deletor):
    try:
        res = db.meeting.update_one(
            {'_id':ObjectId(meeting_id), }, 
            {'$unset':{
                "main_files."+file_id: 1
            }}
            )
        return res.modified_count
    except:
        traceback.print_exc()
    return 0

def op_meeting_rm_ref_file(meeting_id, file_id, deletor):
    try:
        res = db.meeting.update_one(
            {'_id':ObjectId(meeting_id), }, 
            {'$unset':{
                "ref_files."+file_id: 1
            }}
            )
        return res.modified_count
    except:
        traceback.print_exc()
    return 0

def op_meeting_add_record(meeting_id, content, usr):
    id = create_uuid()
    ct = str(time.time())
    try:
        res = db.meeting.update_one(
            {'_id':ObjectId(meeting_id)}, 
            {'$set':{
                "record."+id:[ct, usr, content]
            }}
            )
        return res.modified_count
    except:
        traceback.print_exc()
    return 0

def op_meeting_add_commont(meeting_id, author, content, linkid=''):
    id = create_uuid()
    ct = str(time.time())
    try:
        res = db.meeting.update_one(
            {'_id':ObjectId(meeting_id)}, 
            {'$set':{
                "comments."+id:[ct, author, content, linkid]
            }}
            )
        return res.modified_count
    except:
        traceback.print_exc()
    return 0


def __redirect_url():
    return request.args.get('next') or url_for('index') #request.referrer

def __inner_register():
    username = request.form.get('username', '')
    passwd = request.form.get('passwd', '')
    if username != '' and passwd != '':
        if User.regist_new_usr(username, passwd):
            return True
    return False

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        if __inner_register():
            return redirect(__redirect_url())
    return render_template('register.html')

@app.route('/api/register', methods=['POST'])
def api_register():
    if request.method == 'POST':
        if __inner_register():
            return 'success'
    return 'failed'


def __inner_login():
    username = request.form.get('username', '')
    passwd = request.form.get('passwd', '')
    remember = (request.form.get('remember', 'no') == 'yes')
    if username != '' and passwd != '':
        cUser = User.get_from_name(username)
        if cUser != None and cUser.passwd == passwd:
            if login_user(cUser, remember = remember):
                return True
    return False

@app.route("/login", methods=["GET", "POST"])
def login():
    next = request.args.get('next')
    if request.method == 'POST':
        if __inner_login():
            return redirect(__redirect_url())
        else:
            flash('用户名或密码错误')
    return render_template('login.html', next=next)

@app.route("/api/login", methods=["GET", "POST"])
def api_login():
    if request.method == 'POST':
        if __inner_login():
            return 'success'
    return 'failed'


@app.route('/login_failed')
def login_failed():
    next = __redirect_url()
    if next != None and next.find('api') > -1:
        return 'login_failed'
    return redirect(url_for('login', next=next))

@app.route("/logout")
@login_required
def logout():
    logout_user()
    return "logout"

@app.route('/index')
@login_required
def index():
    rr = get_meeting_list()
    return render_template('index.html', src=rr)


@app.route('/api/listmeeting')
@login_required
def listmeeting():
    rr = get_meeting_list()
    return jsonify(**rr)


@app.route('/meeting_info/<uid>')
@login_required
def meeting_info(uid):
    struid = str(uid)
    cfg = op_meeting_find_by_id(struid)
    if cfg:
        return render_template('meeting_info.html', src=cfg)
    else:
        flash('查询会议失败')
        return redirect(url_for('index'))

def append_timestr_to_file_array(files):
    for key in files:
        try:
            objkey = ObjectId(key)
            t = datetime_to_timetimestr(objkey.generation_time)
            files[key].insert(0, t)
            # files[key].append(t)
        except:
            traceback.print_exc()


@app.route('/api/meeting_info/<uid>')
@login_required
def api_meeting_info(uid):
    struid = str(uid)
    cfg = op_meeting_find_by_id(struid)
    if cfg:
        # return Response(
        #     json_util.dumps(cfg),
        #     mimetype='application/json'
        #     )
        cfg['id'] = struid
        cfg.pop('_id')
        mainfiles = cfg.get('main_files')
        if mainfiles:
            append_timestr_to_file_array(mainfiles)
        reffiles = cfg.get('ref_files')
        if reffiles:
            append_timestr_to_file_array(reffiles)
        return jsonify(cfg)
    else:
        return 'null'

def __inner_create_meeting():
    meeting_name = request.form.get('meeting_name')
    creator = current_user.name
    meeting_name = secure_name(meeting_name)
    id = op_meeting_new(meeting_name, creator)
    return id

@app.route('/create_meeting', methods=['GET', 'POST'])
@login_required
def create_meeting():
    if request.method == 'POST':
        id = __inner_create_meeting()
        if id != None:
            flash('创建成功')
            return redirect(url_for('meeting_info', uid=id))
        else:
            flash('创建失败')
            return redirect(url_for('index'))
    else:
        return render_template('create_meeting.html')

@app.route('/api/create_meeting', methods=['POST'])
@login_required
def api_create_meeting():
    if request.method == 'POST':
        id = __inner_create_meeting()
        if id != None:
            return 'success'
    else:
        return 'failed'


def __inner_delete_meeting():
    mid = request.form.get('meeting_id')
    deleter = current_user.name
    res = op_meeting_delete(mid)
    return res

@app.route('/delete_meeting', methods=['POST'])
@login_required
def delete_meeting():
    if request.method == 'POST':
        res = __inner_delete_meeting()
        if res != None and res == 1:
            flash('删除完毕') #todo
        else:
            flash('删除过程中发生错误')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))
        

@app.route('/api/delete_meeting', methods=['POST'])
@login_required
def api_delete_meeting():
    if request.method == 'POST':
        res = __inner_delete_meeting()
        if res != None and res == 1:
            return 'success'
    return 'failed'

def __inner_upload_usr_file():
    mid = request.form.get('meeting')
    author = current_user.name
    sucess = False
    file = request.files['file']
    if file != None:
        filename = file.filename
        temp = request.form.get('filename')
        if temp != None:
            filename = temp
        filename = secure_filename(filename)
        if file_ext_check(filename):
            cnt = op_meeting_add_main_file(mid, file, filename, author)
            return cnt
    return None

@app.route('/upload_usr_file', methods=['POST'])
@login_required
def upload_usr_file():
    if request.method == 'POST':
        cnt = __inner_upload_usr_file()
        if cnt != None and cnt == 1:
            flash('上传成功')
        else:
            flash('上传过程中发生错误')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))

@app.route('/api/upload_usr_file', methods=['POST'])
@login_required
def api_upload_usr_file():
    if request.method == 'POST':
        cnt = __inner_upload_usr_file()
        if cnt != None and cnt == 1:
            return 'success'
    return 'failed'

def __inner_upload_ref_file():
    mid = request.form.get('meeting')
    author = current_user.name
    file = request.files['file']
    if file != None:
        filename = file.filename
        temp = request.form.get('filename')
        if temp != None:
            filename = temp
        filename = secure_filename(filename)
        if file_ext_check(filename):
            cnt = op_meeting_add_ref_file(mid, file, filename, author)
            return cnt
    return None

@app.route('/upload_ref_file', methods=['POST'])
@login_required
def upload_ref_file():
    if request.method == 'POST':
        cnt = __inner_upload_ref_file()
        if cnt != None and cnt == 1:
            flash('上传成功')
        else:
            flash('上传过程中发生错误')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))
        
@app.route('/api/upload_ref_file', methods=['POST'])
@login_required
def api_upload_ref_file():
    if request.method == 'POST':
        cnt = __inner_upload_ref_file()
        if cnt != None and cnt == 1:
            return 'success'
    return 'failed'



def __inner_delete_usr_file():
    mid = request.form.get('meeting')
    deletor = current_user.name
    file = request.form.get('file')
    if file != None:
        cnt = op_meeting_rm_usr_file(mid, file, deletor)
        return cnt
    return None

@app.route('/delete_usr_file', methods=['POST'])
@login_required
def delete_usr_file():
    if request.method == 'POST':
        cnt = __inner_delete_usr_file()
        if cnt != None and cnt == 1:
            flash('删除成功')
        else:
            flash('删除过程中发生错误')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))

        
@app.route('/api/delete_usr_file', methods=['POST'])
@login_required
def api_delete_usr_file():
    if request.method == 'POST':
        cnt = __inner_delete_usr_file()
        if cnt != None and cnt == 1:
            return 'success'
    return 'failed'

def __inner_delete_ref_file():
    mid = request.form.get('meeting')
    deletor = current_user.name
    file = request.form.get('file')
    if file != None:
        cnt = op_meeting_rm_ref_file(mid, file, deletor)
        return cnt
    return None

@app.route('/delete_ref_file', methods=['POST'])
@login_required
def delete_ref_file():
    if request.method == 'POST':
        cnt = __inner_delete_ref_file()
        if cnt != None and cnt == 1:
            flash('删除成功')
        else:
            flash('删除过程中发生错误')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))

        
@app.route('/api/delete_ref_file', methods=['POST'])
@login_required
def api_delete_ref_file():
    if request.method == 'POST':
        cnt = __inner_delete_ref_file()
        if cnt != None and cnt == 1:
            return 'success'
    return 'failed'



@app.route('/download_file', methods=['GET', 'POST'])
def download_file():
    #mid = ''
    fid = ''
    if request.method == 'POST':
        #mid = request.form.get('meeting_id')
        fid = request.form.get('file_id')
    else:
        #mid = request.args.get('meeting_id')
        fid = request.args.get('file_id')
    if fid != None:
        objid = ObjectId(fid)
        fp = fs.get(objid)
        if fp:
            return send_file(fp, as_attachment=True, attachment_filename=fp.filename, conditional=True)
    return redirect(__redirect_url())
        #a, b = os.path.split(fpath)
        # return send_from_directory(a, b, as_attachment=True)

def __inner_update_record():
    mid = request.form.get('meeting_id')
    data = request.form.get('data')
    usr = current_user.name
    cnt = op_meeting_add_record(mid, data, usr)
    return cnt

@app.route('/update_record', methods=['POST'])
@login_required
def update_record():
    if request.method == 'POST':
        cnt = __inner_update_record()
        if cnt != None and cnt == 1:
            flash('成功记录')
        else:
            flash('记录时发生错误')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))

@app.route('/api/update_record', methods=['GET', 'POST'])
@login_required
def api_update_record():
    if request.method == 'POST':
        cnt = __inner_update_record()
        if cnt != None and cnt == 1:
            return 'success'
    return 'failed'

def __inner_update_comments():
    mid = request.form.get('meeting_id')
    data = request.form.get('data')
    link = request.form.get('link', '')
    usr = current_user.name
    cnt = op_meeting_add_commont(mid, usr, data, link)
    return cnt

@app.route('/update_comments', methods=['GET', 'POST'])
@login_required
def update_comments():
    if request.method == 'POST':
        cnt = __inner_update_comments()
        if cnt != None and cnt == 1:
            flash('发表成功')
        else:
            flash('发表失败')
    if request.referrer:
        return redirect(request.referrer)
    else:
        return redirect(url_for('index'))

@app.route('/api/update_comments', methods=['GET', 'POST'])
@login_required
def api_update_comments():
    if request.method == 'POST':
        cnt = __inner_update_comments()
        if cnt != None and cnt == 1:
            return 'success'
    return 'failed'

@app.route('/')
@login_required
def web_root():
    return redirect(url_for('index'))

@app.route('/about')
def web_about():
    return render_template('about.html')

@app.route('/help')
def web_help():
    return render_template('help.html')

@app.route('/contact')
def web_contact():
    return render_template('contact.html')

@app.route('/share')
def web_share():
    return redirect(url_for('web_about'))

@app.route('/star')
def web_star():
    return redirect(url_for('web_about'))

# @app.route('/listusr')
# def list_usr():
#     res = ''
#     nr = User.get_all_usr_name()
#     for i in nr:
#         res += i
#         res += '\n'
#     return res

def main():
    # import logging
    # from logging.handlers import RotatingFileHandler
    # handler = RotatingFileHandler('sv4app.log', maxBytes=10000, backupCount=1)
    # handler.setLevel(logging.INFO)
    # app.logger.addHandler(handler)
    app.debug = False
    app.run(host='0.0.0.0')

if __name__ == '__main__':
    main()
